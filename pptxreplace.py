"""Powerpoint file replace text module"""
import logging  # noqa: E402


from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER

from config import replace_substring


# def remove_metadata_from_app_xml(prs):
#     """There is currently no functionality for handling app.xml so
#     have to find the part and then alter its blob manually
#     """
#     package_parts = prs.part.package.parts
#     for part in package_parts:
#         if part.partname.endswith('app.xml'):
#             app_xml_part = part
#     app_xml = app_xml_part.blob.decode('utf-8')
#     tags_to_remove = ('Company', 'Manager', 'HyperlinkBase')
#     for tag in tags_to_remove:
#         pattern = f'<{tag}>.*<\/{tag}>'
#         app_xml = re.sub(pattern, '', app_xml)
#     app_xml_part.blob = bytearray(app_xml, 'utf-8')


def process_shape(shape_parent, data, replaced, verbose=False):
    for shape in shape_parent.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for key, value in data.items():
                        txt_old = run.text
                        run.text, n = replace_substring(
                            run.text, value, key)
                        replaced[key] += n
                        print(
                            f'{txt_old} -> {run.text}') if verbose else None

        if shape.has_table:
            table = shape.table
            for cell in table.iter_cells():
                # here you can access the text in cell by using
                # cell.text
                # just remember that the shape object refers to the table in this context not the cell
                for key, value in data.items():
                    txt_old = cell.text
                    cell.text, n = replace_substring(
                        cell.text, value, key)
                    replaced[key] += n
                    print(
                        f'{txt_old} -> {cell.text}') if verbose else None

        if shape.has_chart:
            chart = shape.chart
            for series in chart.series:
                for key, value in data.items():
                    txt_old = series.name
                    txt_new, n = replace_substring(
                        txt_old, value, key)
                    replaced[key] += n
                    series.name.replace(txt_old, txt_new)
                    print(
                        f'{txt_old} -> {txt_new}') if verbose else None
            for categ in chart.plots[0].categories:
                for key, value in data.items():
                    txt_old = categ
                    categ, n = replace_substring(
                        txt_old, value, key)
                    replaced[key] += n
                    print(
                        f'{txt_old} -> {categ}') if verbose else None

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            process_shape(shape, data, replaced, verbose)

        if shape.is_placeholder:
            ph = shape.placeholder_format
            if ph.type == PP_PLACEHOLDER.FOOTER or ph.type == PP_PLACEHOLDER.HEADER:
                # print('%d, %s' % (ph.idx, ph.type))
                tidx = shape_parent.shapes[ph.idx]
                sp = tidx.element
                sp.getparent().remove(shape.element)


def replace_pptx(file_path: str, new_path: str, data: dict) -> None:
    prs = Presentation(file_path)

    # remove_metadata_from_app_xml(prs)

    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    replaced = {value: 0 for value in data.keys()}
    verbose = False
    for slide in prs.slides:
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            for key, value in data.items():
                txt_old = notes_slide.notes_text_frame.text
                notes_slide.notes_text_frame.text, n = replace_substring(
                    notes_slide.notes_text_frame.text, value, key)
                replaced[key] += n
                print(
                    f'{txt_old} -> {notes_slide.notes_text_frame.text}') if verbose else None

        process_shape(slide, data, replaced, verbose)

    logging.info(f'Replacements:\n  {replaced}')

    prs.save(new_path)
    logging.info(f'New PPTX file saved to {new_path}')
